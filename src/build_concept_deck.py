# -*- coding: utf-8 -*-
"""개념 해설 PPT — 기초 용어 → 일반화 → 전이(transfer) → 우리 프로젝트.
비유 없이, 용어 정의 중심, 기술적. 네이비+틸, Pretendard, 16:9, 네이티브 표.
build_deck.py와 동일한 시각 시스템 재사용.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

NAVY=(0x0E,0x1B,0x33); NAVY2=(0x18,0x29,0x4C); INK=(0x18,0x26,0x3D); MUTED=(0x5C,0x6B,0x82)
FAINT=(0x9A,0xA8,0xBE); BLUE=(0x25,0x63,0xEB); TEAL=(0x12,0xB3,0xA6); TEALDK=(0x0C,0x8A,0x80)
AMBER=(0xF1,0x9E,0x0B); GREEN=(0x1F,0xA8,0x5C); RED=(0xD6,0x49,0x49)
BG=(0xF4,0xF7,0xFB); CARD=(0xFF,0xFF,0xFF); LINE=(0xDD,0xE5,0xF0); TEALBG=(0xE7,0xF6,0xF3)
NAVYTX=(0xC7,0xD3,0xE6); WHITE=(0xFF,0xFF,0xFF)
FONT="Pretendard"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def _tf(rPr,name):
    for tag in ("a:latin","a:ea","a:cs"):
        el=rPr.find(qn(tag))
        if el is None: el=etree.SubElement(rPr,qn(tag))
        el.set("typeface",name)
def style(run,r):
    f=run.font; f.size=Pt(r.get("s",16)); f.bold=r.get("b",False); f.italic=r.get("i",False)
    f.color.rgb=RGBColor(*r.get("c",INK)); f.name=r.get("f",FONT)
    rPr=run._r.get_or_add_rPr(); _tf(rPr,r.get("f",FONT))
    if "spc" in r: rPr.set("spc",str(int(r["spc"])))
def tb(s,l,t,w,h,paras,anchor=MSO_ANCHOR.TOP,wrap=True):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=b.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,p in enumerate(paras):
        pa=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pa.alignment=p.get("align",PP_ALIGN.LEFT)
        if "before" in p: pa.space_before=Pt(p["before"])
        pa.space_after=Pt(p.get("after",0))
        if "line" in p: pa.line_spacing=p["line"]
        for r in p["runs"]: rr=pa.add_run(); rr.text=r["t"]; style(rr,r)
    return b
def shp(s,kind,l,t,w,h,fill=None,line=None,lw=1.0,adj=None):
    sh=s.shapes.add_shape(kind,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(*fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=RGBColor(*line); sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    if adj is not None:
        try: sh.adjustments[0]=adj
        except Exception: pass
    return sh
def rrect(s,l,t,w,h,fill=None,line=None,lw=1.0,adj=0.06): return shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h,fill,line,lw,adj)
def oval(s,l,t,w,h,fill=None,line=None,lw=1.0): return shp(s,MSO_SHAPE.OVAL,l,t,w,h,fill,line,lw)
def setbg(s,c): f=s.background.fill; f.solid(); f.fore_color.rgb=RGBColor(*c)
def shptext(sh,paras,anchor=MSO_ANCHOR.MIDDLE):
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=0; tf.margin_bottom=0
    for i,p in enumerate(paras):
        pa=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pa.alignment=p.get("align",PP_ALIGN.CENTER)
        pa.space_after=Pt(p.get("after",0))
        for r in p["runs"]: rr=pa.add_run(); rr.text=r["t"]; style(rr,r)

TOTAL="12"
def header(s,num,title,sub):
    tb(s,0.6,0.34,9,0.34,[{"runs":[{"t":"AI 공격탐지 · 기술 개념 해설","s":11,"b":True,"c":TEALDK,"spc":60}]}])
    tb(s,11.0,0.34,1.733,0.34,[{"runs":[{"t":f"{num} / {TOTAL}","s":10.5,"c":FAINT}],"align":PP_ALIGN.RIGHT}])
    tb(s,0.6,0.78,12.1,0.7,[{"runs":[{"t":num+"  ","s":27,"b":True,"c":TEAL},{"t":title,"s":27,"b":True,"c":NAVY}]}])
    tb(s,0.6,1.5,12.1,0.4,[{"runs":[{"t":sub,"s":13.5,"c":MUTED}]}])

def table(s,x,y,colw,rowh,data,header_fill=NAVY,header_tc=WHITE,fs=11,hl=None):
    hl=hl or set(); yy=y
    for ri,row in enumerate(data):
        xx=x
        for ci,cell in enumerate(row):
            if ri==0:
                rrect(s,xx,yy,colw[ci],rowh,header_fill,adj=0.0); tc=header_tc; bold=True
            else:
                fill=TEALBG if ri in hl else CARD
                shp(s,MSO_SHAPE.RECTANGLE,xx,yy,colw[ci],rowh,fill=fill,line=LINE,lw=0.75); tc=INK; bold=(ci==0)
            al=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            tb(s,xx+0.06,yy,colw[ci]-0.12,rowh,[{"runs":[{"t":str(cell),"s":fs,"b":bold,"c":tc}],"align":al}],MSO_ANCHOR.MIDDLE)
            xx+=colw[ci]
        yy+=rowh

def defrow(s,x,y,w,h,kr,en,body,note=None,accent=TEAL):
    """용어 정의 카드 한 줄: 한글명(영문) + 정의 + (선택)우리 프로젝트 메모."""
    rrect(s,x,y,w,h,CARD,LINE,1,adj=0.045); rrect(s,x,y+0.13,0.09,h-0.26,accent,adj=0.5)
    runs=[{"t":kr,"s":15,"b":True,"c":NAVY}]
    if en: runs.append({"t":f"   {en}","s":11,"b":True,"c":accent})
    tb(s,x+0.32,y+0.15,w-0.5,0.38,[{"runs":runs}])
    paras=[{"runs":[{"t":body,"s":12,"c":INK}],"line":1.13}]
    if note:
        paras.append({"runs":[{"t":"▸ 우리 프로젝트  ","s":10.5,"b":True,"c":TEALDK},{"t":note,"s":10.5,"c":MUTED}],"before":4,"line":1.08})
    tb(s,x+0.32,y+0.58,w-0.52,h-0.7,paras)

CT,CW3,G=2.0,3.78,0.33; LM=0.6  # 3-column geometry

# ===== S1 표지 =====
s=prs.slides.add_slide(BLANK); setbg(s,NAVY)
cx,cy=11.75,2.3
for r,col,lw in [(2.45,NAVY2,1.5),(1.75,BLUE,1.25),(1.1,TEAL,1.5)]: oval(s,cx-r,cy-r,2*r,2*r,None,col,lw)
nd=oval(s,11.42,1.97,0.66,0.66,NAVY2,TEAL,1.75); shptext(nd,[{"runs":[{"t":"AI","s":15,"b":True,"c":WHITE}]}])
pill=rrect(s,0.9,0.92,3.5,0.5,TEAL,adj=0.5); shptext(pill,[{"runs":[{"t":"사이버보안 WE-Meet · 개념 해설","s":12,"b":True,"c":NAVY}]}])
tb(s,0.9,1.95,9.5,1.7,[
   {"runs":[{"t":"개념부터 우리 프로젝트까지","s":40,"b":True,"c":WHITE}],"line":1.05},
   {"runs":[{"t":"용어 정의 중심 ","s":22,"b":True,"c":TEAL},{"t":"기술 해설","s":22,"b":True,"c":NAVYTX}],"before":8}])
tb(s,0.92,3.95,11,0.5,[{"runs":[{"t":"분류·피처·평가 → 일반화 → 전이(transfer) → 우리의 발견과 개선","s":17,"c":NAVYTX}]}])
chips=["기초 용어","평가 지표","일반화·전이","우리 프로젝트"]
lx=0.92
for i,c in enumerate(chips):
    ch=rrect(s,lx,5.0,2.45,0.7,NAVY2,NAVY2,1,adj=0.18); shptext(ch,[{"runs":[{"t":c,"s":14,"b":True,"c":WHITE}]}])
    if i<3: tb(s,lx+2.45,5.0,0.5,0.7,[{"runs":[{"t":"→","s":20,"b":True,"c":TEAL}],"align":PP_ALIGN.CENTER}],MSO_ANCHOR.MIDDLE)
    lx+=2.95
tb(s,0.92,6.7,11,0.4,[{"runs":[{"t":"네트워크 침입 탐지(CSE-CIC-IDS2018)를 예시로  ·  비유 없이 정의 중심으로 설명","s":12.5,"c":FAINT}]}])

# ===== S2 (01) 학습 지도 — 갭을 메우는 계단 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"01","이 자료의 길 — 어디서 어디로 가는가","아는 것(왼쪽)에서 우리 프로젝트의 핵심(오른쪽)까지, 한 칸씩 정의하며 올라간다")
ladder=[("①","분류·모델·학습","02"),("②","피처·라벨","03"),("③","학습/시험 분리\n과적합·지름길","04"),
        ("④","성능 지표\n재현율·임계값","05"),("⑤","일반화·분포이동","06"),("⑥","전이(transfer)","07"),
        ("⑦","전이의 증명\n(oracle)","08"),("⑧","개선: 시간맥락","09"),("⑨","위험점수·근거·설명","10~12")]
# 2행 배치
cols=5; bw=2.3; bh=1.15; gx=0.18; gy=0.5
x0=(13.333-(cols*bw+(cols-1)*gx))/2; y0=2.25
for i,(no,lab,sl) in enumerate(ladder):
    row=i//cols; col=i%cols
    x=x0+col*(bw+gx); y=y0+row*(bh+gy)
    cur=(i<=1)  # 사용자가 대략 아는 지점
    key=(no=="⑥")
    fill=TEAL if key else (CARD)
    rrect(s,x,y,bw,bh,fill,(TEALDK if key else LINE),1.25 if key else 1,adj=0.08)
    tcol=WHITE if key else NAVY
    tb(s,x+0.16,y+0.12,bw-0.3,0.34,[{"runs":[{"t":no+" ","s":15,"b":True,"c":(WHITE if key else TEAL)},{"t":("핵심" if key else ("기초" if cur else "")),"s":10,"b":True,"c":(WHITE if key else FAINT)}]}])
    tb(s,x+0.16,y+0.46,bw-0.3,0.62,[{"runs":[{"t":lab,"s":12.5,"b":True,"c":tcol}],"line":1.02}])
    tb(s,x+bw-0.62,y+bh-0.3,0.55,0.26,[{"runs":[{"t":"p."+sl,"s":9,"c":(NAVYTX if key else FAINT)}],"align":PP_ALIGN.RIGHT}])
tb(s,0.6,5.7,12.1,0.9,[
  {"runs":[{"t":"지금 감 잡으신 곳 ①②","s":12.5,"b":True,"c":MUTED},{"t":"  →  도착점 ⑥ 전이(transfer)","s":12.5,"b":True,"c":TEALDK},{"t":"  : '왜 새 공격을 못 잡나'의 답.","s":12.5,"c":INK}],"after":3,"line":1.2},
  {"runs":[{"t":"③④⑤는 ⑥을 이해하기 위한 디딤돌입니다. 한 장씩 용어를 정의하며 갭을 메웁니다.","s":12,"c":MUTED}],"line":1.2}])

# ===== S3 (02) 용어1: 모델·학습·분류 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"02","기본 용어 (1) — 모델 · 학습 · 분류","머신러닝이 '예측하는 함수'를 데이터로 만드는 과정")
defrow(s,0.6,2.05,12.13,1.18,"모델 (model)","= 함수 f(x)",
  "입력 x(숫자 묶음)를 받아 출력 ŷ(예측)를 내놓는 함수. 내부에 조절 가능한 값(파라미터 θ)이 있어, θ를 바꾸면 예측이 바뀐다.",
  "입력 = 통신기록 1건, 출력 = 정상/공격 예측.")
defrow(s,0.6,3.35,12.13,1.18,"학습 (training)","= θ를 맞추는 과정",
  "정답(라벨)이 붙은 데이터를 많이 주고, 예측이 정답과 어긋난 정도(손실, loss)를 최소화하도록 파라미터 θ를 자동으로 조정한다. 규칙을 사람이 직접 쓰지 않는다.",
  "수십~수백만 건의 (통신기록, 정답) 쌍으로 θ를 맞춤.")
defrow(s,0.6,4.65,12.13,1.18,"이진 분류 (binary classification)","= 출력이 두 범주",
  "출력이 0 또는 1 두 가지인 분류 문제. 우리는 0 = 정상(Benign), 1 = 공격(Attack)으로 둔다.",
  "사용 모델: XGBoost = 결정트리(decision tree) 수백 개를 차례로 더해 오차를 줄이는 앙상블(gradient boosting). 표 형태 데이터에서 딥러닝보다 강함.",accent=BLUE)

# ===== S4 (03) 용어2: 피처·피처벡터·라벨 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"03","기본 용어 (2) — 피처 · 피처 벡터 · 라벨","모델은 원본이 아니라 '숫자로 요약된 입력'만 본다")
defrow(s,0.6,2.05,12.13,1.18,"피처 (feature)","= 입력 변수 하나",
  "한 대상을 수치화한 개별 측정값. 모델은 피처들의 값만 보고 판단한다.",
  "통신기록 1건에서 뽑은 값들 — 패킷 길이, 통신 지속시간, 패킷 간 간격(IAT), 플래그 수 등.")
defrow(s,0.6,3.35,12.13,1.18,"피처 벡터 (feature vector)","= 피처들의 묶음 x",
  "한 샘플을 이루는 피처 값들을 순서대로 나열한 숫자 묶음. x = (x₁, x₂, …, x₇₈) 처럼 길이 78의 벡터.",
  "CICFlowMeter라는 도구가 통신기록 1건당 약 78개 통계를 자동 추출 → 이 78개가 모델 입력.")
defrow(s,0.6,4.65,12.13,1.18,"라벨 (label) · 지도학습 (supervised)","= 정답 y",
  "각 샘플에 붙은 정답값 y. 정답이 있는 (x, y) 쌍으로 배우는 방식을 지도학습이라 한다.",
  "데이터에 '이 흐름은 공격/정상'이 라벨로 붙어 있어 지도학습이 가능. (단, 일부 라벨엔 오류가 알려져 있음)",accent=BLUE)

# ===== S5 (04) 용어3: 분리·과적합·지름길 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"04","기본 용어 (3) — 학습/시험 분리 · 과적합 · 지름길","'외운 점수'와 '진짜 실력'을 가르는 개념들")
defrow(s,0.6,2.0,12.13,1.12,"학습/시험 분리 (train / test split)","",
  "데이터를 학습용과 시험용으로 나눠, 모델이 한 번도 보지 못한 시험용에서 성능을 잰다. 같은 데이터로 평가하면 암기 여부만 재는 셈이라 의미가 없다.",accent=TEAL)
defrow(s,0.6,3.22,12.13,1.12,"과적합 (overfitting)","",
  "학습 데이터에선 성능이 높은데 시험 데이터에선 떨어지는 상태. 데이터 고유의 잡음·우연한 상관까지 외워버려 일반화에 실패한 것.",accent=AMBER)
defrow(s,0.6,4.44,12.13,1.12,"지름길 학습 (shortcut learning)","",
  "진짜 인과(공격 행위)가 아니라, 그 데이터에서 정답과 우연히 강하게 맞물린 피처 하나에 의존하는 것. 다른 환경에선 무너진다.",accent=AMBER)
band=rrect(s,0.6,5.75,12.13,0.95,NAVY,adj=0.1)
tb(s,0.95,5.92,11.6,0.7,[
  {"runs":[{"t":"우리 사례 — ","s":12.5,"b":True,"c":TEAL},{"t":"1일치로 학습하니 F1 = 1.000(100%). 그러나 피처 중요도 1위 'Fwd Seg Size Min' 하나가 ","s":12,"c":NAVYTX},{"t":"99.2%","s":13,"b":True,"c":WHITE},{"t":" 차지 = 전형적 지름길.","s":12,"c":NAVYTX}],"line":1.15}])

# ===== S6 (05) 용어4: 성능 지표 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"05","기본 용어 (4) — 성능 지표","무엇을 '잘했다'고 부를지 정의한다 (오탐은 이미 아는 개념)")
# 좌: 혼동행렬
tb(s,0.6,2.0,5.4,0.34,[{"runs":[{"t":"혼동행렬 (confusion matrix)","s":13.5,"b":True,"c":NAVY}]}])
cm=[["", "예측: 공격", "예측: 정상"],
    ["실제: 공격","TP (적중)","FN (미탐)"],
    ["실제: 정상","FP (오탐)","TN (정상적중)"]]
table(s,0.6,2.45,[1.7,1.95,1.95],0.6,cm,fs=11)
tb(s,0.6,4.45,5.6,2.0,[
  {"runs":[{"t":"FP 오탐","s":11.5,"b":True,"c":RED},{"t":" : 정상을 공격이라 잘못 알림","s":11.5,"c":INK}],"after":5,"line":1.1},
  {"runs":[{"t":"FN 미탐","s":11.5,"b":True,"c":AMBER},{"t":" : 공격을 놓침 (보안에서 더 위험)","s":11.5,"c":INK}],"after":5,"line":1.1},
  {"runs":[{"t":"TP 적중","s":11.5,"b":True,"c":GREEN},{"t":" : 공격을 공격으로 맞힘","s":11.5,"c":INK}],"line":1.1}])
# 우: 지표 정의
rx=6.55
metrics=[("재현율 (recall)","= TP / (TP+FN)","실제 공격 중 잡아낸 비율. 우리 핵심 지표.",TEALDK),
         ("정밀도 (precision)","= TP / (TP+FP)","공격이라 한 것 중 진짜 공격 비율.",MUTED),
         ("임계값 (threshold τ)","점수→판정 경계","모델은 '공격일 점수 s'를 낸다. s ≥ τ면 공격으로 선언. τ를 낮추면 재현율↑·오탐↑ (줄다리기).",BLUE),
         ("FP / 10만","운영 오탐 단위","정상 10만 건당 오탐 수. '오탐 예산'을 고정하고 재현율을 비교해야 공정.",MUTED),
         ("AP (Average Precision)","순위 품질","임계값 전 구간의 평균 성능. 높아도 특정 τ에서 재현율 0일 수 있어 함께 본다.",MUTED)]
yy=2.0
for kr,ex,body,clr in metrics:
    rrect(s,rx,yy,6.18,0.86,CARD,LINE,1,adj=0.07); rrect(s,rx,yy+0.12,0.08,0.62,clr,adj=0.5)
    tb(s,rx+0.28,yy+0.1,5.85,0.34,[{"runs":[{"t":kr+"  ","s":12.5,"b":True,"c":NAVY},{"t":ex,"s":10.5,"b":True,"c":clr}]}])
    tb(s,rx+0.28,yy+0.45,5.85,0.38,[{"runs":[{"t":body,"s":10.5,"c":MUTED}],"line":1.05}])
    yy+=0.95

# ===== S7 (06) 용어5: 일반화·분포이동 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"06","기본 용어 (5) — 일반화 · 분포 이동","전이를 이해하기 위한 마지막 디딤돌")
defrow(s,0.6,2.05,12.13,1.15,"일반화 (generalization)","",
  "학습에 없던 새 데이터에서의 성능. 시험용에서의 점수가 곧 일반화 성능이다. 머신러닝이 진짜 노리는 목표.",accent=TEAL)
defrow(s,0.6,3.3,12.13,1.15,"i.i.d. 가정 (독립 동일분포)","",
  "학습 데이터와 시험 데이터가 '같은 분포에서 독립적으로 뽑혔다'는 전제. 데이터를 무작위로 학습/시험 분리하면 이 전제가 성립해 점수가 높게 나온다.",accent=BLUE)
defrow(s,0.6,4.55,12.13,1.15,"분포 이동 (distribution shift)","",
  "시험 데이터의 분포가 학습 때와 달라지는 것. i.i.d. 가정이 깨져 일반화가 무너진다. 실제 운영(시간이 흐르며 환경·공격이 변함)에서 흔히 발생.",accent=AMBER)
band=rrect(s,0.6,5.85,12.13,0.85,NAVY,adj=0.11)
tb(s,0.95,6.0,11.6,0.6,[{"runs":[
  {"t":"표기 정리 — ","s":12.5,"b":True,"c":TEAL},
  {"t":"입력 X, 라벨 Y의 관계는  P(X, Y) = P(X) · P(Y | X).  ","s":12,"c":NAVYTX},
  {"t":"다음 장에서 이 P(Y|X)가 핵심.","s":12,"b":True,"c":WHITE}],"line":1.1}])

# ===== S8 (07) ★ 전이(transfer) 정의 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"07","전이 (transfer) — 이 프로젝트의 핵심 개념","한 분포에서 배운 판정 규칙이 다른 분포로 '옮겨가는가'의 문제")
defrow(s,0.6,1.95,12.13,1.0,"전이 (transfer)","",
  "원천 분포(과거 날짜)에서 학습한 결정 함수를, 목표 분포(미래 날짜·새 공격)에 적용하는 것. 잘 옮겨가지 못하면 '전이 실패' → 새 공격을 못 잡는다.",accent=TEAL)
# 두 종류의 분포이동 대비
bx=0.6; bw=5.96
b1=rrect(s,bx,3.12,bw,2.05,CARD,LINE,1,adj=0.045); rrect(s,bx,3.28,0.09,1.73,GREEN,adj=0.5)
tb(s,bx+0.34,3.3,bw-0.6,0.4,[{"runs":[{"t":"공변량 이동 (covariate shift)","s":14,"b":True,"c":NAVY}]}])
tb(s,bx+0.34,3.78,bw-0.62,1.3,[
  {"runs":[{"t":"P(X)는 변하지만  ","s":12,"c":INK},{"t":"P(Y|X)는 그대로","s":12.5,"b":True,"c":GREEN}],"after":5,"line":1.15},
  {"runs":[{"t":"입력 분포만 바뀌고, '이 입력이면 이 정답'이라는 관계는 유지. 비교적 다루기 쉬움.","s":11.5,"c":MUTED}],"line":1.15}])
bx2=6.77
b2=rrect(s,bx2,3.12,bw,2.05,CARD,LINE,1,adj=0.045); rrect(s,bx2,3.28,0.09,1.73,RED,adj=0.5)
tb(s,bx2+0.34,3.3,bw-0.6,0.4,[{"runs":[{"t":"개념(조건부) 이동 (concept shift)","s":14,"b":True,"c":NAVY}]}])
tb(s,bx2+0.34,3.78,bw-0.62,1.3,[
  {"runs":[{"t":"P(Y|X) 자체가 변함","s":12.5,"b":True,"c":RED},{"t":"   ← 우리 경우","s":11,"b":True,"c":RED}],"after":5,"line":1.15},
  {"runs":[{"t":"'이 입력이면 공격'이라는 관계가 달라짐. 과거에 배운 규칙이 미래엔 안 맞음. 어렵다.","s":11.5,"c":MUTED}],"line":1.15}])
band=rrect(s,0.6,5.4,12.13,1.3,NAVY,adj=0.07)
tb(s,0.95,5.56,11.6,1.05,[
  {"runs":[{"t":"왜 우리가 '개념 이동'인가  ","s":12.5,"b":True,"c":TEAL}],"after":3},
  {"runs":[{"t":"날짜마다 등장하는 공격 종류가 거의 겹치지 않는다 → '이 피처값이면 공격'이라는 관계가 날짜마다 달라진다. 따라서 과거 날짜로 배운 P(Y|X)가 미래 날짜(다른 공격)에는 들어맞지 않는다. ","s":12,"c":NAVYTX},
           {"t":"병목은 모델 용량·피처 부족이 아니라 P(Y|X)의 시간적 변화(전이 실패).","s":12,"b":True,"c":WHITE}],"line":1.2}])

# ===== S9 (08) 전이의 증명 — oracle =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"08","전이임을 어떻게 증명했나 — 오라클 비교","'정보가 없어서'가 아니라 '안 옮겨가서'임을 수치로 분리")
tb(s,0.6,2.0,6.0,2.4,[
  {"runs":[{"t":"날짜 교차 평가 ","s":12.5,"b":True,"c":NAVY},{"t":"(Leave-One-Group-Out by day)","s":10.5,"c":MUTED}],"after":3,"line":1.1},
  {"runs":[{"t":"한 날짜를 통째로 시험셋으로 빼고 나머지 날짜로 학습 → 실제 운영(과거→미래)과 같은 조건.","s":11.5,"c":INK}],"after":10,"line":1.15},
  {"runs":[{"t":"오라클 (target-day oracle)","s":12.5,"b":True,"c":NAVY}],"after":3,"line":1.1},
  {"runs":[{"t":"'만약 그날 공격이 학습에 있었다면?'을 가정 — 그날 데이터로 학습해 그날을 시험. 정보의 상한선.","s":11.5,"c":INK}],"after":10,"line":1.15},
  {"runs":[{"t":"raw-kNN·Deep SAD(표현학습)도 동일하게 실패","s":11.5,"b":True,"c":TEALDK},{"t":" → '피처에 정보가 없어서'가 아님.","s":11.5,"c":INK}],"line":1.15}])
hd=["공격","오라클\n(정보 상한)","운영\n(실제)","해석"]
rows=[hd,
 ["Bot(봇넷)","0.996","≈ 0","전이 실패"],
 ["DoS-GE","1.00","0.03","전이 실패"],
 ["Web","0.68","≈ 0","전이 실패"]]
table(s,6.7,2.05,[1.5,1.7,1.3,1.5],0.66,rows,fs=10.5,hl={1})
tb(s,6.7,5.05,6.0,1.3,[
  {"runs":[{"t":"결론  ","s":12.5,"b":True,"c":TEALDK},{"t":"오라클은 높은데 운영은 0 → 피처 안에 공격을 가를 정보는 충분하다. 못 잡는 이유는 그 정보가 과거→미래로 전이되지 않기 때문.","s":12,"c":INK}],"line":1.2}])

# ===== S10 (09) 개선 — 시간맥락(temporal) =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"09","개선 — 시간맥락(temporal) 피처","'시점 간 행동 패턴'을 더해 전이로 막힌 부분을 일부 복원")
tb(s,0.6,2.0,5.7,0.34,[{"runs":[{"t":"왜 / 무엇을 (IP 없이)","s":13.5,"b":True,"c":NAVY}]}])
tb(s,0.6,2.45,5.75,3.0,[
  {"runs":[{"t":"문제: ","s":12,"b":True,"c":AMBER},{"t":"흐름(flow) 1건은 시점 독립이라 '과거 맥락'이 없다. 봇넷의 주기적 반복 연결이 한 건엔 안 드러남.","s":11.5,"c":INK}],"after":8,"line":1.18},
  {"runs":[{"t":"추가 피처: ","s":12,"b":True,"c":TEALDK},{"t":"같은 목적지 포트를 키로, 과거~현재 시간창(1·10·60·300·600초)에서 —","s":11.5,"c":INK}],"after":4,"line":1.18},
  {"runs":[{"t":"· 연결 시도 수","s":11.5,"c":INK}],"after":2,"line":1.1},
  {"runs":[{"t":"· SYN/RST·연결 실패 비율","s":11.5,"c":INK}],"after":2,"line":1.1},
  {"runs":[{"t":"· 연결 간격의 규칙성(주기성)","s":11.5,"c":INK}],"after":8,"line":1.1},
  {"runs":[{"t":"causal: ","s":12,"b":True,"c":TEALDK},{"t":"미래 정보는 쓰지 않음(데이터 누수 방지).","s":11.5,"c":INK}],"line":1.15}])
hd=["공격","추가 전","추가 후"]
rows=[hd,["Bot(봇넷)","0.009","0.49"],["DoS","0.51","0.62"],["DDoS","0.38","0.50"],["BruteForce","0.74","0.98"],["Web","0.02","≈ 0"]]
table(s,6.7,2.05,[2.4,1.8,1.83],0.56,rows,fs=11.5,hl={1})
tb(s,6.7,5.55,6.0,1.1,[
  {"runs":[{"t":"봇넷 0.009 → 0.49.","s":12.5,"b":True,"c":TEALDK},{"t":" 주기적 반복 연결이 시간창 통계에 드러나 크게 개선. Web은 요청 내용(payload)이 데이터에 없어 미해결.","s":12,"c":INK}],"line":1.2}])

# ===== S11 (10) 위험점수·보정 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"10","위험점수와 보정 (calibration · conformal)","탐지 점수를 '믿을 수 있는 우선순위'로 바꾸는 단계")
defrow(s,0.6,2.0,12.13,1.05,"보정 (calibration)","",
  "모델이 낸 점수를 실제 확률에 맞추는 것. '공격 확률 0.8'이라면 그런 사례 100건 중 80건이 실제 공격이어야 함. 어긋난 정도는 ECE로 측정.",accent=BLUE)
defrow(s,0.6,3.15,12.13,1.05,"문제 — 단순 보정은 분포이동에서 깨짐","",
  "sigmoid 보정을 적용했더니 날짜가 바뀌자 확률이 크게 어긋났다(ECE 0.13~0.25). 분포이동 때문에 '확률처럼 보이는 값'을 믿을 수 없게 됨.",accent=AMBER)
defrow(s,0.6,4.3,12.13,1.05,"conformal p-value","",
  "정상(benign) 점수 분포에서 이 샘플이 얼마나 드문 꼬리에 있는지를 순위로 계산한 값. 분포를 가정하지 않아 분포이동에 강하다.",accent=TEAL)
defrow(s,0.6,5.45,12.13,1.05,"Neyman-Pearson 임계값","",
  "오탐율 상한(예: FP ≤ 18/10만)을 제약으로 두고 그 안에서 재현율을 최대화하도록 임계값을 정함. → 위험점수 0~100은 이 p-value 기반 운영 점수.",accent=TEAL)

# ===== S12 (11) 근거·설명 SHAP/LLM =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"11","근거(SHAP)와 설명(LLM)의 역할","판정은 ML이, 설명은 LLM이 — 권한을 분리한다")
defrow(s,0.6,2.05,12.13,1.35,"SHAP (SHapley Additive exPlanations)","근거",
  "협조 게임이론의 Shapley 값을 적용해, 각 피처가 이 예측을 '기준값에서 얼마나 밀었는지' 기여도를 계산한다. 기여도의 합 = (예측값 − 기준값). 어떤 피처 때문에 공격으로 봤는지 제시.",
  "TreeExplainer(트리 모델 전용, 정확·고속)로 이벤트별 상위 기여 피처를 추출 → 설명 입력으로 사용.",accent=BLUE)
defrow(s,0.6,3.55,12.13,1.35,"LLM (대규모 언어모델)","설명만",
  "탐지·점수가 만든 '변경 불가(immutable) 근거 JSON'을 받아 사람이 읽을 문장으로만 바꾼다. 판정·점수·우선순위는 절대 바꾸지 않는다.",
  "이유 — 재현성 확보 + 프롬프트 인젝션·환각이 판정에 끼어드는 것을 차단. LLM이 틀려도 탐지 결과는 불변.",accent=TEAL)
band=rrect(s,0.6,5.25,12.13,1.45,NAVY,adj=0.07)
tb(s,0.95,5.42,11.6,0.4,[{"runs":[{"t":"역할 분리 원칙","s":13,"b":True,"c":TEAL}]}])
tb(s,0.95,5.82,11.6,0.8,[
  {"runs":[{"t":"탐지·위험점수·우선순위 = ML이 결정(재현 가능).  ","s":12.5,"c":NAVYTX},{"t":"SHAP = 왜 그렇게 봤는지 근거.  ","s":12.5,"c":NAVYTX},{"t":"LLM = 그 근거를 말로 풀어주는 통역.","s":12.5,"b":True,"c":WHITE}],"line":1.2}])

# ===== S13 (12) 전체 구조 + 핵심·한계 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"12","전체 구조와 핵심 메시지","개념들이 하나의 파이프라인으로 연결된다")
steps=["데이터\n(78 피처+시간맥락)","XGBoost\n탐지(공격/정상)","conformal\n위험점수 0~100","SHAP\n근거 피처","LLM\n자연어 설명","대시보드\n위험순 화면"]
SW,SG=1.86,0.27; SLM=(13.333-(6*SW+5*SG))/2; SCT,SCH=2.15,1.25
for i,lab in enumerate(steps):
    x=SLM+i*(SW+SG); rrect(s,x,SCT,SW,SCH,CARD,LINE,1,adj=0.08)
    c=oval(s,x+SW/2-0.24,SCT+0.16,0.48,0.48,NAVY,TEAL,1.5); shptext(c,[{"runs":[{"t":str(i+1),"s":13,"b":True,"c":WHITE}]}])
    tb(s,x+0.05,SCT+0.72,SW-0.1,0.5,[{"runs":[{"t":lab,"s":10,"b":True,"c":NAVY}],"align":PP_ALIGN.CENTER,"line":1.0}])
    if i<5: tb(s,x+SW,SCT,SG,SCH,[{"runs":[{"t":"→","s":15,"b":True,"c":TEAL}],"align":PP_ALIGN.CENTER}],MSO_ANCHOR.MIDDLE)
# 핵심/한계 2칸
k=rrect(s,0.6,3.8,5.96,2.6,CARD,LINE,1,adj=0.04); rrect(s,0.6,3.96,0.09,2.28,TEAL,adj=0.5)
tb(s,0.95,3.98,5.4,0.4,[{"runs":[{"t":"핵심 메시지","s":15,"b":True,"c":NAVY}]}])
tb(s,0.95,4.5,5.5,1.8,[
  {"runs":[{"t":"•  ","s":12,"b":True,"c":TEAL},{"t":"못 잡는 원인은 모델·표현 부족이 아니라 전이(P(Y|X)의 시간 변화).","s":12,"c":INK}],"after":7,"line":1.15},
  {"runs":[{"t":"•  ","s":12,"b":True,"c":TEAL},{"t":"시간맥락 피처로 미관측 공격을 0 → 0.5급까지 부분 개선.","s":12,"c":INK}],"after":7,"line":1.15},
  {"runs":[{"t":"•  ","s":12,"b":True,"c":TEAL},{"t":"평가는 무작위 분할이 아닌 날짜 교차로만 보고.","s":12,"c":INK}],"line":1.15}])
l=rrect(s,6.77,3.8,5.96,2.6,CARD,LINE,1,adj=0.04); rrect(s,6.77,3.96,0.09,2.28,AMBER,adj=0.5)
tb(s,7.12,3.98,5.4,0.4,[{"runs":[{"t":"남은 한계","s":15,"b":True,"c":NAVY}]}])
tb(s,7.12,4.5,5.5,1.8,[
  {"runs":[{"t":"•  ","s":12,"b":True,"c":AMBER},{"t":"미래에 처음 나올 공격은 미리 학습 불가 — 전이의 본질적 한계.","s":12,"c":INK}],"after":7,"line":1.15},
  {"runs":[{"t":"•  ","s":12,"b":True,"c":AMBER},{"t":"Web 공격은 요청 내용(payload)이 데이터에 없어 미탐.","s":12,"c":INK}],"after":7,"line":1.15},
  {"runs":[{"t":"•  ","s":12,"b":True,"c":AMBER},{"t":"단일 모델은 운영 수준 미달 → 다층 방어(규칙+네트워크+단말+로그) 필요.","s":12,"c":INK}],"line":1.15}])

base=r"C:\Users\WannaGoHome\Desktop\내 문서\coss\사이버보안 WE-MEET\output\발표자료"
os.makedirs(base,exist_ok=True)
out=os.path.join(base,"사이버보안 WE-Meet 개념·전이 기술해설.pptx")
try:
    prs.save(out)
except PermissionError:
    out=os.path.join(base,"사이버보안 WE-Meet 개념·전이 기술해설_v2.pptx")
    prs.save(out)
print("저장:",out,"| 슬라이드:",len(prs.slides._sldIdLst))
