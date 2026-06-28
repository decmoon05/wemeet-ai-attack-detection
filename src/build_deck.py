# -*- coding: utf-8 -*-
"""프로젝트 소개용 PPT (결과 반영본) — 네이비+틴트, Pretendard, 16:9.
문제→접근→정직한 평가→temporal 성과→다층방어→일정. python-pptx 네이티브.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

NAVY=(0x0E,0x1B,0x33); NAVY2=(0x18,0x29,0x4C); INK=(0x18,0x26,0x3D); MUTED=(0x5C,0x6B,0x82)
FAINT=(0x9A,0xA8,0xBE); BLUE=(0x25,0x63,0xEB); TEAL=(0x12,0xB3,0xA6); TEALDK=(0x0C,0x8A,0x80)
AMBER=(0xF1,0x9E,0x0B); GREEN=(0x1F,0xA8,0x5C); RED=(0xD6,0x49,0x49)
BG=(0xF4,0xF7,0xFB); CARD=(0xFF,0xFF,0xFF); LINE=(0xDD,0xE5,0xF0); TEALBG=(0xE7,0xF6,0xF3)
NAVYTX=(0xC7,0xD3,0xE6); WHITE=(0xFF,0xFF,0xFF)
FONT="Pretendard"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def _tf(rPr,name):
    for tag in ("a:latin","a:ea","a:cs"):
        el=rPr.find(qn(tag));
        if el is None: el=etree.SubElement(rPr,qn(tag))
        el.set("typeface",name)
def style(run,r):
    f=run.font; f.size=Pt(r.get("s",16)); f.bold=r.get("b",False); f.italic=r.get("i",False)
    f.color.rgb=RGBColor(*r.get("c",INK)); f.name=r.get("f",FONT)
    rPr=run._r.get_or_add_rPr(); _tf(rPr,r.get("f",FONT))
    if "spc" in r: rPr.set("spc",str(int(r["spc"])))
def tb(s,l,t,w,h,paras,anchor=MSO_ANCHOR.TOP,wrap=True):
    b=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=b.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,p in enumerate(paras):
        pa=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pa.alignment=p.get("align",PP_ALIGN.LEFT)
        if "before" in p: pa.space_before=Pt(p["before"])
        pa.space_after=Pt(p.get("after",0))
        if "line" in p: pa.line_spacing=p["line"]
        for r in p["runs"]: rr=pa.add_run(); rr.text=r["t"]; style(rr,r)
    return b
def shp(s,kind,l,t,w,h,fill=None,line=None,lw=1.0,adj=None):
    sh=s.shapes.add_shape(kind,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(*fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=RGBColor(*line); sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    if adj is not None:
        try: sh.adjustments[0]=adj
        except Exception: pass
    return sh
def rrect(s,l,t,w,h,fill=None,line=None,lw=1.0,adj=0.06): return shp(s,MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h,fill,line,lw,adj)
def oval(s,l,t,w,h,fill=None,line=None,lw=1.0): return shp(s,MSO_SHAPE.OVAL,l,t,w,h,fill,line,lw)
def setbg(s,c): f=s.background.fill; f.solid(); f.fore_color.rgb=RGBColor(*c)
def shptext(sh,paras,anchor=MSO_ANCHOR.MIDDLE):
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08); tf.margin_top=0; tf.margin_bottom=0
    for i,p in enumerate(paras):
        pa=tf.paragraphs[0] if i==0 else tf.add_paragraph(); pa.alignment=p.get("align",PP_ALIGN.CENTER)
        pa.space_after=Pt(p.get("after",0))
        for r in p["runs"]: rr=pa.add_run(); rr.text=r["t"]; style(rr,r)

TOTAL="09"
def header(s,num,title,sub):
    tb(s,0.6,0.34,8,0.34,[{"runs":[{"t":"사이버보안 WE-Meet","s":11,"b":True,"c":TEALDK,"spc":60}]}])
    tb(s,11.0,0.34,1.733,0.34,[{"runs":[{"t":f"{num} / {TOTAL}","s":10.5,"c":FAINT}],"align":PP_ALIGN.RIGHT}])
    tb(s,0.6,0.78,11.6,0.7,[{"runs":[{"t":num+"  ","s":28,"b":True,"c":TEAL},{"t":title,"s":28,"b":True,"c":NAVY}]}])
    tb(s,0.6,1.5,11.6,0.4,[{"runs":[{"t":sub,"s":13.5,"c":MUTED}]}])

def table(s,x,y,colw,rowh,data,header_fill=NAVY,header_tc=WHITE,fs=11,hl=None):
    """간단 표: data[0]=헤더. hl=강조할 행 인덱스 set(배경 틴트)."""
    hl=hl or set()
    yy=y
    for ri,row in enumerate(data):
        xx=x
        for ci,cell in enumerate(row):
            if ri==0:
                rrect(s,xx,yy,colw[ci],rowh,header_fill,adj=0.0); tc=header_tc; bold=True
            else:
                fill=TEALBG if ri in hl else CARD
                shp(s,MSO_SHAPE.RECTANGLE,xx,yy,colw[ci],rowh,fill=fill,line=LINE,lw=0.75); tc=INK; bold=(ci==0)
            al=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            tb(s,xx+0.06,yy,colw[ci]-0.12,rowh,[{"runs":[{"t":str(cell),"s":fs,"b":bold,"c":tc}],"align":al}],MSO_ANCHOR.MIDDLE)
            xx+=colw[ci]
        yy+=rowh

# ===== S1 표지 =====
s=prs.slides.add_slide(BLANK); setbg(s,NAVY)
cx,cy=11.75,2.3
for r,col,lw in [(2.45,NAVY2,1.5),(1.75,BLUE,1.25),(1.1,TEAL,1.5)]: oval(s,cx-r,cy-r,2*r,2*r,None,col,lw)
nd=oval(s,11.42,1.97,0.66,0.66,NAVY2,TEAL,1.75); shptext(nd,[{"runs":[{"t":"AI","s":15,"b":True,"c":WHITE}]}])
pill=rrect(s,0.9,0.92,2.85,0.5,TEAL,adj=0.5); shptext(pill,[{"runs":[{"t":"사이버보안 WE-Meet","s":12.5,"b":True,"c":NAVY}]}])
tb(s,0.9,1.95,9.2,1.1,[{"runs":[{"t":"AI","s":44,"b":True,"c":TEAL},{"t":" 기반 공격탐지 에이전트","s":44,"b":True,"c":WHITE}]}])
tb(s,0.92,3.0,11,0.5,[{"runs":[{"t":"CSE-CIC-IDS2018 네트워크 흐름 기반 침입 탐지 · 중간 진행 보고","s":18,"c":NAVYTX}]}])
tb(s,0.92,4.7,8,0.32,[{"runs":[{"t":"프로젝트 소개 자료 · 소집 멘토링 대체 제출본","s":12,"b":True,"c":TEAL,"spc":40}]}])
chips=["문제·데이터","구성·방법","평가·결과","한계·다음"]
lx=0.92
for i,c in enumerate(chips):
    ch=rrect(s,lx,5.1,2.45,0.7,NAVY2,NAVY2,1,adj=0.18); shptext(ch,[{"runs":[{"t":c,"s":14,"b":True,"c":WHITE}]}])
    if i<3: tb(s,lx+2.45,5.1,0.5,0.7,[{"runs":[{"t":"→","s":20,"b":True,"c":TEAL}],"align":PP_ALIGN.CENTER}],MSO_ANCHOR.MIDDLE)
    lx+=2.95
tb(s,0.92,6.8,11,0.4,[{"runs":[{"t":"팀  팀장 · 팀원      |      컴퓨터융합학부      |      2026 하기 계절학기","s":12.5,"c":FAINT}]}])

# ===== S2 왜 필요한가 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"01","문제 정의와 목표","규칙 기반 탐지의 한계를 보완하되, 성능을 신뢰할 수 있게 측정")
cards=[("배경",["포트 스캔·DoS/DDoS·봇넷 등 공격 유형이 다양","규칙 기반 탐지는 알려진 패턴만 탐지(신규·변형에 취약)","트래픽·로그가 방대해 사람이 직접 보기 어려움"]),
       ("해결할 점",["공개 벤치마크의 높은 정확도가 실제 성능을 과대평가","학습에 없던 새 공격(미관측)의 탐지가 특히 어려움","어느 공격을 어디까지 탐지 가능한지 수치로 제시 필요"]),
       ("목표",["정상/공격 탐지 + 위험점수 + 근거 + 설명을 통합","무작위가 아닌 날짜 교차로 실제 성능 측정","위험순으로 분석가가 우선 처리할 화면 제공"])]
LM,CT,CH,G=0.6,2.02,3.0,0.33; CW=(12.133-2*G)/3
for i,(h,items) in enumerate(cards):
    x=LM+i*(CW+G); rrect(s,x,CT,CW,CH,CARD,LINE,1,adj=0.05); rrect(s,x,CT+0.16,0.09,CH-0.32,TEAL,adj=0.5)
    tb(s,x+0.34,CT+0.28,CW-0.5,0.45,[{"runs":[{"t":h,"s":18,"b":True,"c":NAVY}]}])
    tb(s,x+0.34,CT+0.95,CW-0.62,CH-1.1,[{"runs":[{"t":"•  ","s":13,"c":TEAL,"b":True},{"t":it,"s":13,"c":INK}],"after":7,"line":1.1} for it in items])

# ===== S_DATA 데이터셋 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"02","데이터셋 — CSE-CIC-IDS2018","캐나다 CIC 제작, AWS에 모의 기업망 구축 후 실제 공격을 수행·기록한 공개 벤치마크")
# 좌: 개요 카드
rrect(s,0.6,2.0,3.7,3.5,CARD,LINE,1,adj=0.04)
tb(s,0.9,2.2,3.2,0.4,[{"runs":[{"t":"규모·구성","s":15,"b":True,"c":NAVY}]}])
tb(s,0.9,2.75,3.2,2.6,[
  {"runs":[{"t":"총 흐름(flow)  ","s":12,"c":MUTED},{"t":"약 666만 건","s":12.5,"b":True,"c":INK}],"after":6},
  {"runs":[{"t":"수집 기간  ","s":12,"c":MUTED},{"t":"10일 (2/14~3/2)","s":12.5,"b":True,"c":INK}],"after":6},
  {"runs":[{"t":"흐름 통계  ","s":12,"c":MUTED},{"t":"80여 개 (CICFlowMeter)","s":12.5,"b":True,"c":INK}],"after":6},
  {"runs":[{"t":"공격 유형  ","s":12,"c":MUTED},{"t":"14종 / 6계열","s":12.5,"b":True,"c":INK}],"after":6},
  {"runs":[{"t":"라벨  ","s":12,"c":MUTED},{"t":"정상(Benign)·공격","s":12.5,"b":True,"c":INK}],"after":10},
  {"runs":[{"t":"※ 정제판은 IP·시각 제거 →","s":10.5,"c":MUTED}],"after":1,"line":1.05},
  {"runs":[{"t":"   원본 CSV의 Timestamp로 시간맥락 별도 생성","s":10.5,"c":MUTED}],"line":1.05},
])
# 우: 일자별 공격 표
hd=["일자","공격 계열","공격 흐름 수"]
rows=[hd,["2/14","BruteForce(FTP·SSH)","9.4만"],["2/15·16","DoS(GoldenEye·Hulk 등)","20.5만"],
      ["2/20·21","DDoS(LOIC·HOIC)","77.7만"],["2/22·23","Web(BF·XSS·SQLi)","0.09만"],
      ["2/28·3/1","Infiltration(라벨불신)","11.9만"],["3/2","Bot(봇넷)","14.5만"]]
table(s,4.6,2.0,[1.5,4.2,2.0],0.5,rows,fs=11.5,hl={6})
tb(s,4.6,5.7,8,0.4,[{"runs":[{"t":"※ 봇넷(Bot, 3/2)은 학습에 전혀 없는 '미관측 공격'으로 두어 일반화를 시험","s":12,"c":TEALDK,"b":True}]}])

# ===== S3 접근(5요소) =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"03","시스템 구성 — 5단계 파이프라인","입력 흐름을 탐지·점수화하고 근거와 설명을 분석가에게 제공")
steps=[("데이터","CSE-CIC-IDS2018\n80여 흐름 통계","정상/공격 라벨"),
       ("탐지","RandomForest /\nXGBoost","정상·공격 이진"),
       ("위험점수","0~100점 환산","conformal 보정"),
       ("근거","SHAP 기여도","상위 피처 추출"),
       ("설명·화면","LLM 자연어 설명","Streamlit 대시보드")]
SW,SG=2.0,0.53; SLM=(13.333-(5*SW+4*SG))/2; SCT,SCH=2.5,2.2
for i,(lab,l1,l2) in enumerate(steps):
    x=SLM+i*(SW+SG); rrect(s,x,SCT,SW,SCH,CARD,LINE,1,adj=0.06)
    c=oval(s,x+SW/2-0.3,SCT+0.2,0.6,0.6,NAVY,TEAL,1.75); shptext(c,[{"runs":[{"t":str(i+1),"s":15,"b":True,"c":WHITE}]}])
    tb(s,x,SCT+0.9,SW,0.34,[{"runs":[{"t":lab,"s":14,"b":True,"c":NAVY}],"align":PP_ALIGN.CENTER}])
    tb(s,x,SCT+1.28,SW,0.7,[{"runs":[{"t":l1,"s":10.5,"c":INK}],"align":PP_ALIGN.CENTER,"line":1.05},
                            {"runs":[{"t":l2,"s":10,"c":MUTED}],"align":PP_ALIGN.CENTER,"before":2}])
    if i<4: tb(s,x+SW,SCT,SG,SCH,[{"runs":[{"t":"→","s":18,"b":True,"c":TEAL}],"align":PP_ALIGN.CENTER}],MSO_ANCHOR.MIDDLE)
band=rrect(s,0.85,5.25,11.633,1.45,NAVY,adj=0.07)
tb(s,1.2,5.45,11,0.4,[{"runs":[{"t":"역할 분리 원칙","s":13,"b":True,"c":TEAL}]}])
tb(s,1.2,5.85,11,0.7,[{"runs":[{"t":"• 탐지와 위험점수는 RandomForest/XGBoost가 결정한다(재현 가능, 수치 안정).","s":12.5,"c":NAVYTX}],"after":3},
                      {"runs":[{"t":"• LLM은 판정을 바꾸지 않고, 탐지 결과(피처 근거)를 사람이 읽을 문장으로만 변환한다 — 환각·프롬프트 인젝션의 영향을 차단.","s":12.5,"c":NAVYTX}]}])

# ===== S4 정직한 평가(핵심 발견) =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"04","평가 방법과 그 과정에서의 발견","평가 방식에 따라 성능이 크게 달라지는 것을 확인하고 평가 체계를 재설계")
# 좌: 대비 박스
bx=rrect(s,0.6,2.05,5.6,3.35,CARD,LINE,1,adj=0.04)
tb(s,0.95,2.28,5.0,0.4,[{"runs":[{"t":"평가 방식별 결과 (이진 탐지)","s":15,"b":True,"c":NAVY}]}])
tb(s,0.95,2.85,5.1,0.45,[{"runs":[{"t":"① 무작위 분할","s":13,"b":True,"c":MUTED},{"t":"  (데이터를 무작위로 학습/시험 분리)","s":11,"c":FAINT}]}])
tb(s,0.95,3.28,5.1,0.5,[{"runs":[{"t":"F1 0.99","s":20,"b":True,"c":RED},{"t":"  실제 배포 성능과 무관한 과대평가","s":12.5,"c":INK}]}])
tb(s,0.95,4.15,5.1,0.45,[{"runs":[{"t":"② 날짜 교차","s":13,"b":True,"c":TEALDK},{"t":"  (과거 일자로 학습 → 미래 일자로 시험)","s":11,"c":FAINT}]}])
tb(s,0.95,4.58,5.1,0.5,[{"runs":[{"t":"대부분 0","s":20,"b":True,"c":AMBER},{"t":"  학습에 없던 공격은 거의 못 탐지","s":12.5,"c":INK}]}])
# 우: 발견 3
rx=6.5
findings=[("지름길(shortcut) 피처","단일 피처(Fwd Seg Size Min)가 모델 중요도의 99%를 차지. 공격 행위가 아니라 데이터 생성 특성을 학습한 것."),
          ("AP가 높아도 실제 탐지율은 0","순위 지표(AP)는 높게 나와도, 운영 임계값을 적용하면 탐지되는 공격이 거의 없는 경우가 다수."),
          ("성능 저하의 원인은 전이(transfer)","같은 날짜 안에서는 99.6%까지 구분 가능. 과거→미래로 갈 때 데이터 분포가 변하는 것이 원인.")]
yy=2.05
for h,d in findings:
    rrect(s,rx,yy,6.0,1.02,CARD,LINE,1,adj=0.06); rrect(s,rx,yy+0.14,0.09,0.74,AMBER,adj=0.5)
    tb(s,rx+0.32,yy+0.13,5.55,0.36,[{"runs":[{"t":h,"s":13.5,"b":True,"c":NAVY}]}])
    tb(s,rx+0.32,yy+0.5,5.55,0.46,[{"runs":[{"t":d,"s":11,"c":MUTED}],"line":1.05}])
    yy+=1.12
tb(s,0.6,5.6,12,0.5,[{"runs":[{"t":"무작위 분할의 높은 수치는 사용하지 않고, 날짜 교차로 측정한 값을 실제 성능으로 보고한다.","s":13,"b":True,"c":TEALDK}],"align":PP_ALIGN.CENTER}])

# ===== S_FOLDS 9-fold 전체 표 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"05","공격 계열별 날짜 교차 평가 (이진 탐지)","과거 일자로 학습 → 각 공격일을 시험. 무작위 분할은 별도(상한선)로만 참고")
hd=["테스트 일자 / 공격","노출","공격 흐름","AP","탐지율\n@100FP","실제\nFP/10만"]
rows=[hd,
 ["2/15 DoS(GoldenEye·Slow)","미관측","5.1만","0.46","0.03","46"],
 ["2/16 DoS(Hulk)","유사 학습","14.5만","0.17","0.55","91,763 ⚠"],
 ["2/20 DDoS(LOIC-HTTP)","미관측","57.5만","0.98","0.00","11"],
 ["2/21 DDoS(HOIC)","유사 학습","20.1만","1.00","0.83","7"],
 ["2/22 Web(BF·XSS·SQLi)","미관측","341","0.00","0.01","96"],
 ["2/23 Web","유사 학습","541","0.35","0.43","100"],
 ["3/2 Bot(봇넷)","미관측","14.5만","0.31","0.00","77"],
]
table(s,0.6,1.95,[3.5,1.5,1.4,1.3,1.7,1.7],0.47,rows,fs=11,hl={4})
tb(s,0.6,5.55,12.2,0.9,[
 {"runs":[{"t":"무작위 분할 시 F1 0.99","s":12.5,"b":True,"c":RED},{"t":" → 날짜 교차에서는 DDoS-HOIC(0.83)만 안정 탐지.","s":12,"c":INK}],"after":3,"line":1.15},
 {"runs":[{"t":"DoS-Hulk는 탐지율 높아도 오탐 폭발(FP 9만)로 사용 불가. Web은 공격 흐름이 수백 건뿐이라 학습 부족.","s":12,"c":INK}],"line":1.15}])

# ===== S5 결과(temporal 성과) =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"06","시간맥락 피처 추가 후 탐지율 변화","흐름별 통계에 시간 흐름 정보를 더해 미관측 공격 탐지를 개선")
# 좌측 설명
tb(s,0.6,2.0,3.0,0.34,[{"runs":[{"t":"추가한 피처 (IP 불필요)","s":12,"b":True,"c":NAVY}]}])
tb(s,0.6,2.4,3.0,3.0,[
   {"runs":[{"t":"포트별 1·10·60·300·600초 ","s":11,"c":INK},{"t":"창의 연결 수","s":11,"c":MUTED}],"after":5,"line":1.1},
   {"runs":[{"t":"SYN/RST·연결 실패 ","s":11,"c":INK},{"t":"비율","s":11,"c":MUTED}],"after":5,"line":1.1},
   {"runs":[{"t":"반복 연결의 ","s":11,"c":INK},{"t":"시간 규칙성","s":11,"c":MUTED}],"after":10,"line":1.1},
   {"runs":[{"t":"임계값은 과거 데이터 기준","s":10.5,"b":True,"c":TEALDK}],"after":2},
   {"runs":[{"t":"(오탐 약 15~18 / 10만 흐름)","s":10,"c":MUTED}]},
])
# 막대 비교 (before→after)
rows=[("Bot(봇넷)",0.009,0.49),("DoS",0.51,0.62),("DDoS",0.38,0.50),("BruteForce",0.74,0.98),("Web",0.02,0.0)]
ry=2.35; rh=0.6; barx=5.3; barw=5.0
tb(s,3.95,2.0,1.3,0.3,[{"runs":[{"t":"공격 유형","s":10.5,"b":True,"c":MUTED}]}])
tb(s,barx,2.0,barw+1.5,0.3,[{"runs":[{"t":"탐지율  ▏회색=추가 전, 청록=추가 후","s":10.5,"c":MUTED}]}])
for lab,b,a in rows:
    tb(s,3.95,ry-0.02,1.4,0.5,[{"runs":[{"t":lab,"s":12,"b":True,"c":NAVY}]}],MSO_ANCHOR.MIDDLE)
    shp(s,MSO_SHAPE.RECTANGLE,barx,ry+0.05,barw,0.15,fill=(0xE3,0xE9,0xF2))
    if b>0: shp(s,MSO_SHAPE.RECTANGLE,barx,ry+0.05,barw*b,0.15,fill=FAINT)
    if a>0: shp(s,MSO_SHAPE.RECTANGLE,barx,ry+0.24,barw*a,0.15,fill=TEAL)
    else: tb(s,barx,ry+0.2,3,0.3,[{"runs":[{"t":"미해결","s":10,"c":RED}]}])
    tb(s,barx+barw+0.15,ry,1.7,0.5,[{"runs":[{"t":f"{b:.2f} → ","s":11,"c":FAINT},{"t":f"{a:.2f}","s":14,"b":True,"c":TEALDK}]}],MSO_ANCHOR.MIDDLE)
    ry+=rh
band=rrect(s,0.85,5.55,11.633,1.1,NAVY,adj=0.08)
tb(s,1.15,5.72,11,0.4,[{"runs":[{"t":"봇넷(Bot) 탐지율 0.009 → 0.49","s":14,"b":True,"c":TEAL}]}])
tb(s,1.15,6.12,11,0.45,[{"runs":[{"t":"봇넷의 주기적 반복 연결이 시간창 통계에 드러나 탐지가 크게 개선됨. Web(XSS/SQLi)은 흐름 통계에 요청 내용(payload)이 없어 현재 구조로는 미탐.","s":12,"c":NAVYTX}],"line":1.1}])

# ===== S_EXTRA 운영검증 + 공격유형분류 + 로그 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"07","추가 검증 — 운영 임계값·공격유형 분류·로그 분석","신뢰성을 위한 세 가지 보강 확인")
# 좌: 운영 임계값 (오라클 vs 과거기반)
rrect(s,0.6,2.0,5.85,2.0,CARD,LINE,1,adj=0.05)
tb(s,0.9,2.18,5.4,0.36,[{"runs":[{"t":"① 운영 임계값 검증 (과거 데이터 기준)","s":13.5,"b":True,"c":NAVY}]}])
table(s,0.9,2.6,[1.5,1.9,1.85],0.34,
   [["공격","과거기준 탐지율","실제 FP/10만"],["Bot","0.41","18"],["DoS","0.62","16"],["BruteForce","0.98","17"]],fs=10.5)
# 우: 공격유형 분류 multiclass
rrect(s,6.7,2.0,6.03,2.0,CARD,LINE,1,adj=0.05)
tb(s,7.0,2.18,5.6,0.36,[{"runs":[{"t":"② 공격 유형 분류(6계열) macro-F1","s":13.5,"b":True,"c":NAVY}]}])
tb(s,7.0,2.7,5.5,0.5,[{"runs":[{"t":"무작위 분할  ","s":12,"c":MUTED},{"t":"0.945","s":18,"b":True,"c":RED}]}])
tb(s,7.0,3.25,5.5,0.5,[{"runs":[{"t":"시간순 분할  ","s":12,"c":MUTED},{"t":"0.219","s":18,"b":True,"c":AMBER}]}])
tb(s,7.0,3.7,5.6,0.3,[{"runs":[{"t":"→ 이진 탐지와 동일: 무작위 평가는 과대평가","s":11,"c":TEALDK}]}])
# 하: 로그 분석
rrect(s,0.6,4.2,12.13,1.9,CARD,LINE,1,adj=0.04)
tb(s,0.9,4.38,11.5,0.36,[{"runs":[{"t":"③ 보안로그 분석 (다층 방어의 로그 레이어)","s":13.5,"b":True,"c":NAVY}]}])
tb(s,0.9,4.85,11.6,1.1,[
  {"runs":[{"t":"봇넷 감염일 Windows 이벤트로그(.evtx) 분석 — 호스트 3대, 이벤트 1만여 건. ","s":12,"c":INK},
           {"t":"서비스 설치(EventID 7045) / 변경(7040) 다수 탐지","s":12,"b":True,"c":TEALDK},
           {"t":" = 봇넷이 재부팅 후에도 살아남기 위한 '지속성(persistence)' 흔적.","s":12,"c":INK}],"after":4,"line":1.15},
  {"runs":[{"t":"네트워크 탐지(NDR)가 못 보는 호스트 내부 활동을 보완. (단, 공격 라벨 부재로 분류가 아닌 이상징후 분석 수준)","s":11.5,"c":MUTED}],"line":1.1}])

# ===== S6 다층 방어 통찰 =====
s=prs.slides.add_slide(BLANK); setbg(s,BG)
header(s,"08","단일 모델의 한계와 다층 방어에서의 위치","본 프로젝트가 전체 보안 체계에서 담당하는 범위")
cards=[("본 프로젝트 범위",["네트워크 흐름 기반 탐지(NDR) 1개 층","보안로그(Windows 이벤트) 분석 보조","정직한 평가·위험점수·근거·설명"],TEAL),
       ("일반적 보안 운영 구성",["SIEM(로그 상관) + NDR(네트워크) + EDR(단말)","규칙 기반 탐지가 주력, ML은 보조","최종 판단·대응은 분석가가 수행"],BLUE),
       ("본 프로젝트의 의의",["단일 ML의 일반화 한계를 수치로 측정","미관측 공격 탐지율 0 → 0.49로 개선","다층 방어가 필요한 이유를 실증"],GREEN)]
for i,(h,items,clr) in enumerate(cards):
    x=LM+i*(CW+G); rrect(s,x,CT,CW,CH,CARD,LINE,1,adj=0.05); rrect(s,x,CT+0.16,0.09,CH-0.32,clr,adj=0.5)
    tb(s,x+0.34,CT+0.28,CW-0.5,0.45,[{"runs":[{"t":h,"s":16,"b":True,"c":NAVY}]}])
    tb(s,x+0.34,CT+0.95,CW-0.62,CH-1.1,[{"runs":[{"t":"•  ","s":12.5,"c":clr,"b":True},{"t":it,"s":12.5,"c":INK}],"after":7,"line":1.12} for it in items])
tb(s,0.6,5.5,12,0.6,[{"runs":[{"t":"단일 모델로 모든 공격을 탐지하는 것은 어려우며, 규칙·네트워크·단말·로그를 겹친 다층 구성과 분석가의 검토가 함께 필요하다.","s":13.5,"c":NAVY}],"align":PP_ALIGN.CENTER,"line":1.15}])

# ===== S7 일정·다음 =====
s=prs.slides.add_slide(BLANK); setbg(s,NAVY)
cx,cy=11.9,5.0
for r,col,lw in [(2.4,NAVY2,1.5),(1.6,BLUE,1.25),(0.95,TEAL,1.5)]: oval(s,cx-r,cy-r,2*r,2*r,None,col,lw)
pill=rrect(s,0.95,1.4,2.85,0.5,TEAL,adj=0.5); shptext(pill,[{"runs":[{"t":"사이버보안 WE-Meet","s":12.5,"b":True,"c":NAVY}]}])
tb(s,11.0,0.34,1.733,0.34,[{"runs":[{"t":"09 / 09","s":10.5,"c":FAINT}],"align":PP_ALIGN.RIGHT}])
tb(s,0.95,2.5,10,0.8,[{"runs":[{"t":"진행 현황과 다음 계획","s":34,"b":True,"c":WHITE}]}])
items=[("현재까지","탐지·날짜교차 평가·시간맥락 개선·SHAP 근거·LLM 설명·대시보드·로그 분석까지 구현 완료"),
       ("주요 결과","날짜 교차 기준 미관측 공격 탐지율 — 봇넷 0.49·DoS 0.62·DDoS 0.50·BruteForce 0.98 (Web 미탐)"),
       ("다음 계획","교수님 피드백 반영 / (여건 시) 원본 데이터 재처리로 호스트 단위 피처 추가 / 결과보고서 정리")]
yy=3.5
for h,d in items:
    rrect(s,0.95,yy,7.3,0.85,None,TEAL,1.25,adj=0.16)
    tb(s,1.2,yy+0.12,1.3,0.6,[{"runs":[{"t":h,"s":13,"b":True,"c":TEAL}]}],MSO_ANCHOR.MIDDLE)
    tb(s,2.55,yy+0.1,5.55,0.66,[{"runs":[{"t":d,"s":11.5,"c":NAVYTX}],"line":1.08}],MSO_ANCHOR.MIDDLE)
    yy+=0.98
tb(s,0.95,6.75,11,0.4,[{"runs":[{"t":"팀 팀장 · 팀원   |   교수님 피드백 부탁드립니다   |   2026 하기 계절학기","s":12.5,"c":FAINT}]}])

import glob as _g
base=r"C:\Users\WannaGoHome\Desktop\내 문서\coss\사이버보안 WE-MEET\output\발표자료"
out=os.path.join(base,"사이버보안 WE-Meet 팀과제1 PPT.pptx")
try:
    prs.save(out)
except PermissionError:
    out=os.path.join(base,"사이버보안 WE-Meet 팀과제1 PPT_v2.pptx")
    prs.save(out)
print("저장:",out)
